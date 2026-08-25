"""Extract text from lecture PDFs and PPTX for quiz generation."""
import json
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation

FILES = [
    ("Lecture 1", r"C:\Users\l\Downloads\DATA COMM LECTURE 1.pdf"),
    ("Lecture 2", r"C:\Users\l\Downloads\LECTURE 2.pdf"),
    ("Lecture 3", r"C:\Users\l\Documents\DATA COMM LECTURE 3 NEW.pptx"),
    ("Lecture 4", r"C:\Users\l\Downloads\DATA COMM LECTURE 4.pdf"),
]

def extract_pdf(path):
    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        pages.append({"page": i + 1, "text": text.strip()})
    return pages

def extract_pptx(path):
    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        pages.append({"page": i + 1, "text": "\n".join(parts)})
    return pages

def main():
    all_content = {}
    for name, path in FILES:
        p = Path(path)
        if not p.exists():
            print(f"MISSING: {path}")
            continue
        if path.lower().endswith(".pdf"):
            pages = extract_pdf(path)
        else:
            pages = extract_pptx(path)
        all_content[name] = pages
    out = Path(__file__).parent / "lecture_content.json"
    out.write_text(json.dumps(all_content, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"Saved {len(all_content)} lectures to {out}")
    for name, pages in all_content.items():
        total_chars = sum(len(p["text"]) for p in pages)
        print(f"  {name}: {len(pages)} pages, {total_chars} chars")

if __name__ == "__main__":
    import sys
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    main()
